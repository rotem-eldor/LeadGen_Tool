"""
Game Lead Finder — Management Presentation Builder
Run: py build_deck.py
Output: Game_Lead_Finder_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette — CrazyLabs light (Amazon Strategy style) ────────────────────────
NAVY        = RGBColor(0xFF, 0xFF, 0xFF)   # slide bg  → white
NAVY_CARD   = RGBColor(0xE8, 0xF7, 0xF9)  # card bg   → light teal
BLUE        = RGBColor(0x00, 0xB4, 0xC8)  # accent    → CL teal
BLUE_LIGHT  = RGBColor(0x00, 0x97, 0xA7)  # light accent → darker teal label
GREEN       = RGBColor(0x2E, 0x7D, 0x32)  # green     → darker for light bg
GREEN_LIGHT = RGBColor(0x43, 0xA0, 0x47)  # green light
AMBER       = RGBColor(0xE6, 0x5C, 0x00)  # amber     → deeper orange on white
RED         = RGBColor(0xC6, 0x28, 0x28)  # red       → deeper red on white
WHITE       = RGBColor(0x1C, 0x3A, 0x3F)  # "white" text → now dark teal body
OFF_WHITE   = RGBColor(0x1C, 0x3A, 0x3F)  # body text → dark teal
MUTED       = RGBColor(0x64, 0x74, 0x8B)  # muted     → slate grey
CARD_LIGHT  = RGBColor(0xB2, 0xE8, 0xEE)  # card light → teal border/divider

W = Inches(13.33)   # WIDE layout width
H = Inches(7.5)     # WIDE layout height

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 (freeform, but addShape uses shape_id)
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return txb

def multiline(slide, lines, x, y, w, h, size=16, color=WHITE,
              bold=False, align=PP_ALIGN.LEFT, line_bold=None):
    """lines = list of str or (str, bold_override)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, lb = line
        else:
            text, lb = line, bold
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = lb
        run.font.name = "Calibri"
    return txb

def card(slide, x, y, w, h, color=None):
    c = color or NAVY_CARD
    r = rect(slide, x, y, w, h, c)
    return r

def accent_bar(slide, x, y, h, color=BLUE):
    rect(slide, x, y, 0.06, h, color)

def slide_number(slide, n, total):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3,
        size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def section_label(slide, label, color=BLUE_LIGHT):
    txt(slide, label.upper(), 0.5, 0.22, 5, 0.3,
        size=10, color=color, bold=True)

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank(prs); bg(s, NAVY)
    # Left accent strip
    rect(s, 0, 0, 0.35, 7.5, BLUE)
    # Big title
    txt(s, "Game Lead Finder", 0.8, 1.6, 9, 1.2,
        size=54, bold=True, color=WHITE)
    # Subtitle
    txt(s, "Automated M&A Intelligence for Mobile Games",
        0.8, 2.9, 9, 0.6, size=22, color=BLUE_LIGHT)
    # Divider
    rect(s, 0.8, 3.65, 4.5, 0.04, BLUE)
    # Meta
    txt(s, "CrazyLabs  ·  Business Development  ·  June 2026",
        0.8, 3.85, 9, 0.4, size=14, color=MUTED)
    # Right decoration — stat callout
    card(s, 9.8, 1.8, 2.8, 2.6, NAVY_CARD)
    accent_bar(s, 9.8, 1.8, 2.6, BLUE)
    txt(s, "265", 10.0, 2.0, 2.4, 0.9, size=52, bold=True,
        color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, "active leads tracked", 10.0, 2.95, 2.4, 0.4,
        size=13, color=MUTED, align=PP_ALIGN.CENTER)
    txt(s, "$60K – $150K/mo", 10.0, 3.4, 2.4, 0.35,
        size=13, color=GREEN_LIGHT, bold=True, align=PP_ALIGN.CENTER)

def slide_intro(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "01 · Introduction")
    txt(s, "What is this tool?", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    bullets = [
        ("🎯", "Purpose", "Identify mobile games worth acquiring before competitors do — automatically."),
        ("📊", "Data source", "Sensor Tower — top 10,000 games by downloads, Google Play + App Store, worldwide."),
        ("⚙️", "How it works", "Every game is scored using a profit model, filtered by range, and sorted by signal strength."),
        ("🔁", "Output", "A live HTML dashboard: 4 tabs, sortable columns, notes, review tracking — rebuilt every run."),
    ]
    for i, (icon, label, desc) in enumerate(bullets):
        bx = 0.5 + i * 3.2
        card(s, bx, 1.6, 3.0, 3.6, NAVY_CARD)
        accent_bar(s, bx, 1.6, 3.6, BLUE)
        txt(s, icon,  bx+0.2, 1.75, 0.6, 0.6, size=28)
        txt(s, label, bx+0.2, 2.45, 2.6, 0.4, size=15, bold=True, color=WHITE)
        txt(s, desc,  bx+0.2, 2.95, 2.6, 2.0, size=12, color=OFF_WHITE)

    # Bottom note
    txt(s, "Two modes:  Small targets  $60K–$150K/mo   ·   Standard M&A  $150K–$2M/mo",
        0.5, 5.45, 12.3, 0.4, size=13, color=MUTED)
    slide_number(s, n, total)

def slide_model(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "02 · Profit Model")
    txt(s, "How we estimate profit", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    # First: RPD explanation
    card(s, 0.5, 1.55, 12.3, 0.95, CARD_LIGHT)
    txt(s, "Revenue Per Download (RPD)  =  90-day Revenue ÷ 90-day Downloads  ÷ 3  =  monthly estimate",
        0.7, 1.7, 11.8, 0.5, size=14, color=OFF_WHITE)
    txt(s, "RPD determines which model applies:", 0.7, 2.05, 8, 0.35, size=12, color=MUTED)

    models = [
        (RED,        "IAA Model",    "RPD < $0.10",
         "Downloads/mo × $0.10 × 50%",
         "Pure ad-monetised games. Downloads drive revenue.\n⚠️  ~74% median error — always verify manually."),
        (BLUE,       "Blend Model",  "RPD ≥ $0.10  &  Rev ≤ $1M/mo",
         "65% × Revenue  +  35% × IAA estimate",
         "Hybrid games with both ads and in-app purchases.\n~33% median error — reasonable for screening."),
        (GREEN,      "IAP High",     "Revenue > $1M/mo",
         "42% × Revenue",
         "Strong IAP titles. Revenue is the dominant signal.\n~30% median error — most reliable model."),
    ]
    for i, (color, title, condition, formula, note) in enumerate(models):
        bx = 0.5 + i * 4.12
        card(s, bx, 2.7, 3.9, 4.1, NAVY_CARD)
        accent_bar(s, bx, 2.7, 4.1, color)
        txt(s, title,     bx+0.2, 2.85, 3.5, 0.45, size=17, bold=True, color=WHITE)
        rect(s, bx+0.2, 3.35, 3.5, 0.04, color)
        txt(s, condition, bx+0.2, 3.45, 3.5, 0.35, size=11, color=color, bold=True)
        txt(s, formula,   bx+0.2, 3.85, 3.5, 0.55, size=12, color=OFF_WHITE, italic=True)
        txt(s, note,      bx+0.2, 4.5,  3.5, 2.1,  size=11, color=MUTED)

    slide_number(s, n, total)

def slide_anchors(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "03 · Calibration Anchors")
    txt(s, "How we validated the model", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    # Stats row
    stats = [("30", "anchor games"), ("2025", "full-year data"), ("3", "model types tested"), ("42%", "IAP margin (calibrated)")]
    for i, (num, label) in enumerate(stats):
        bx = 0.5 + i * 3.1
        card(s, bx, 1.55, 2.8, 1.3, CARD_LIGHT)
        txt(s, num,   bx+0.15, 1.65, 2.5, 0.65, size=36, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
        txt(s, label, bx+0.15, 2.28, 2.5, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # Sample anchors table
    txt(s, "Sample anchor games (known 2025 monthly profit vs. model prediction):",
        0.5, 3.1, 12, 0.35, size=12, color=MUTED)

    headers = ["Game", "Model type", "Actual profit", "Predicted", "Error"]
    col_w   = [3.6, 2.0, 1.9, 1.9, 1.6]
    col_x   = [0.5, 4.1, 6.1, 8.0, 9.9]
    row_h   = 0.42
    row_y_start = 3.55
    rows = [
        ("Zombie Waves",        "IAP High", "$500K", "$1,213K", "+143%"),
        ("SuitU",               "IAP High", "$820K",   "$559K",  "-32%"),
        ("Lamar – Idle Vlogger","IAA",       "$400K",   "$210K",  "-47%"),
        ("Coffee Mania",        "Blend",     "$270K",   "$247K",   "-8% ✓"),
        ("Heroes vs Hordes",    "Blend",     "$150K",   "$179K",  "+19% ✓"),
    ]
    # Header row
    rect(s, 0.5, row_y_start, 11.1, row_h, BLUE)
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        txt(s, hdr, cx+0.1, row_y_start+0.07, cw, row_h-0.1,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Data rows
    for i, row in enumerate(rows):
        ry = row_y_start + row_h + i * row_h
        bg_color = NAVY_CARD if i % 2 == 0 else CARD_LIGHT
        rect(s, 0.5, ry, 11.1, row_h, bg_color)
        err_color = GREEN if "✓" in row[4] else (RED if "-" not in row[4] or int(row[4].replace("%","").replace("+","").replace(" ✓","").strip()) > 30 else AMBER)
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            c = err_color if j == 4 else OFF_WHITE
            txt(s, cell, cx+0.1, ry+0.07, cw, row_h-0.1,
                size=12, color=c, align=PP_ALIGN.LEFT)

    txt(s, "Full calibration: 2 missing games (Dazzly, Royal Escape) — not in 2025 Sensor Tower export",
        0.5, 6.9, 12, 0.35, size=11, color=MUTED, italic=True)
    slide_number(s, n, total)

def slide_exclusions(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "04 · Publisher Exclusions")
    txt(s, "Who we filter out — and why", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    cols = [
        (BLUE, "Hardcoded list  (80+ publishers)",
         "Major publishers we would never acquire — EA, Supercell, King, Tencent, Zynga, Kabam, Playtika, and 70+ more. Manually curated and updated as needed."),
        (AMBER, "Auto-exclude rule",
         "Any publisher appearing 8 or more times in the filtered dataset is excluded automatically. Catches established studios not on the hardcoded list."),
        (GREEN, "Calibration anchors  (30 games)",
         "Known benchmark games used to validate our profit model. Excluded from output so they never appear as new acquisition targets."),
    ]
    for i, (color, title, desc) in enumerate(cols):
        bx = 0.5 + i * 4.12
        card(s, bx, 1.6, 3.9, 4.8, NAVY_CARD)
        accent_bar(s, bx, 1.6, 4.8, color)
        txt(s, title, bx+0.2, 1.78, 3.5, 0.6, size=16, bold=True, color=WHITE)
        rect(s, bx+0.2, 2.45, 3.5, 0.04, color)
        txt(s, desc,  bx+0.2, 2.6,  3.5, 2.8, size=13, color=OFF_WHITE)

    txt(s, "Result: output contains only independent, acquirable studios in the relevant profit range.",
        0.5, 6.65, 12.3, 0.4, size=13, color=MUTED, italic=True)
    slide_number(s, n, total)

def slide_tabs(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "05 · Dashboard Tabs")
    txt(s, "Growing · Stable · Declining · Returning",
        0.5, 0.55, 12, 0.7, size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    tabs = [
        (GREEN,      "📈 Growing",   "> +10% revenue trend",
         "Games gaining momentum — highest priority for outreach.",
         "⚠️  May already exist in Salesforce.\nCross-check before contacting."),
        (BLUE_LIGHT, "➡️  Stable",    "-10% to +10% revenue trend",
         "Consistent performers. Verified NOT in Salesforce — cleanest leads list.",
         "✅  Safe to act on immediately."),
        (RED,        "📉 Declining",  "< -10% revenue trend",
         "Revenue falling. Could be a motivated seller — or a trap. Verify before pursuing.",
         "⚠️  May already exist in Salesforce.\nCross-check before contacting."),
        (AMBER,      "🔄 Returning",  "Cross-platform or prior-run signal",
         "Games found on both Google Play AND App Store, or reappeared from a prior run.",
         "Strongest signal: performing across platforms."),
    ]
    for i, (color, title, condition, desc, note) in enumerate(tabs):
        bx = 0.5 + i * 3.1
        card(s, bx, 1.6, 2.9, 5.1, NAVY_CARD)
        rect(s, bx, 1.6, 2.9, 0.45, color)
        txt(s, title,     bx+0.15, 1.65, 2.6, 0.35, size=15, bold=True, color=NAVY)
        txt(s, condition, bx+0.15, 2.15, 2.6, 0.4,  size=11, color=color, bold=True)
        rect(s, bx+0.15, 2.58, 2.6, 0.03, CARD_LIGHT)
        txt(s, desc,      bx+0.15, 2.7,  2.6, 1.6,  size=12, color=OFF_WHITE)
        rect(s, bx+0.15, 4.35, 2.6, 0.03, CARD_LIGHT)
        txt(s, note,      bx+0.15, 4.45, 2.6, 1.1,  size=11, color=MUTED, italic=True)

    # Current stats bar
    card(s, 0.5, 6.85, 12.3, 0.55, CARD_LIGHT)
    txt(s, "Current run:   62 Growing   ·   66 Stable   ·   38 Declining   ·   97 Returning   =   263 total leads",
        0.8, 6.92, 11.5, 0.38, size=13, color=OFF_WHITE, align=PP_ALIGN.CENTER)
    slide_number(s, n, total)

def slide_automation(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "06 · Automation  —  Coming Soon")
    txt(s, "Execute & Upload", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    # Two cards side by side
    for i, (color, btn, status, title, points) in enumerate([
        (GREEN, "⚡  EXECUTE", "Sensor Tower API — in progress", "One-click full refresh",
         ["Connect to Sensor Tower API directly",
          "Pull latest 90-day data automatically",
          "Re-run analysis and rebuild dashboard",
          "No CSV exports, no manual steps",
          "Schedule: runs every 1st of the month"]),
        (BLUE, "⬆  UPLOAD", "Live today via local server", "Manual CSV upload",
         ["Upload one or multiple CSV files",
          "Drag & drop or file picker",
          "Supports Google Play + App Store together",
          "Cross-platform dedup runs automatically",
          "Dashboard reloads when complete"]),
    ]):
        bx = 0.5 + i * 6.3
        card(s, bx, 1.6, 5.9, 5.2, NAVY_CARD)
        accent_bar(s, bx, 1.6, 5.2, color)
        # Button mockup
        rect(s, bx+0.25, 1.8, 2.2, 0.55, color)
        txt(s, btn, bx+0.25, 1.85, 2.2, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, status, bx+2.6, 1.9, 3.0, 0.4, size=11, color=color, italic=True)
        txt(s, title,  bx+0.25, 2.55, 5.4, 0.45, size=16, bold=True, color=WHITE)
        rect(s, bx+0.25, 3.07, 5.3, 0.03, CARD_LIGHT)
        for k, point in enumerate(points):
            txt(s, "–  " + point, bx+0.3, 3.2 + k*0.5, 5.2, 0.4, size=13, color=OFF_WHITE)

    slide_number(s, n, total)

def slide_salesforce(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "07 · Salesforce Integration  —  Coming Soon")
    txt(s, "Turning leads into a deal workflow", 0.5, 0.55, 12, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    today_items = [
        "Stable tab = verified not in Salesforce",
        "Growing / Declining / Returning may overlap with existing CRM records",
        "Manual cross-check required before outreach",
    ]
    soon_items = [
        ("Real-time dedup",      "Every game checked against Salesforce on each run — no duplicates ever"),
        ("Deal stage overlay",   "See 'In Outreach', 'Meeting Booked', 'Pass' next to each game in the dashboard"),
        ("Auto-push leads",      "New games pushed to Salesforce as Opportunities with pre-filled fields"),
        ("Smart routing",        "Games > $100K auto-assigned to senior managers; geography-based routing"),
    ]

    # Today card (left) — same height as right card
    card(s, 0.5, 1.6, 5.9, 5.4, NAVY_CARD)
    accent_bar(s, 0.5, 1.6, 5.4, AMBER)
    txt(s, "Today", 0.7, 1.75, 5.4, 0.4, size=15, bold=True, color=AMBER)
    rect(s, 0.7, 2.22, 5.1, 0.03, CARD_LIGHT)
    for k, item in enumerate(today_items):
        txt(s, "–  " + item, 0.7, 2.38 + k*0.75, 5.1, 0.65, size=13, color=OFF_WHITE)
    # Bottom callout inside today card
    card(s, 0.75, 5.4, 5.4, 1.35, CARD_LIGHT)
    txt(s, "💡  Once Salesforce is connected, all tabs will automatically show CRM deal status — no more manual cross-checking.",
        0.95, 5.5, 5.0, 1.1, size=12, color=OFF_WHITE, italic=True)

    # Coming soon card (right) — 4 items with comfortable spacing
    card(s, 6.6, 1.6, 6.2, 5.4, NAVY_CARD)
    accent_bar(s, 6.6, 1.6, 5.4, BLUE)
    txt(s, "Coming soon", 6.8, 1.75, 5.7, 0.4, size=15, bold=True, color=BLUE_LIGHT)
    rect(s, 6.8, 2.22, 5.7, 0.03, CARD_LIGHT)
    for k, (label, desc) in enumerate(soon_items):
        txt(s, label, 6.8, 2.38 + k*1.1, 2.2, 0.4, size=13, bold=True, color=WHITE)
        txt(s, desc,  6.8, 2.80 + k*1.1, 5.7, 0.55, size=12, color=MUTED)

    slide_number(s, n, total)

def slide_extras(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    section_label(s, "08 · What Else You Should Know")
    txt(s, "Details that matter", 0.5, 0.55, 10, 0.7,
        size=34, bold=True, color=WHITE)
    rect(s, 0.5, 1.38, 12.3, 0.04, CARD_LIGHT)

    items = [
        (BLUE_LIGHT,  "Two profit ranges",
         "Small targets: $60K–$150K/mo  |  Standard M&A: $150K–$2M/mo\nSame model, different filter. Run separately as needed."),
        (GREEN,       "Cross-platform signal",
         "Games in the Returning tab were found on BOTH Google Play (worldwide) AND App Store (US).\nStrongest possible signal — they perform across ecosystems."),
        (AMBER,       "IAA flag  ⚠️",
         "All IAA-model games are marked with ⚠️ in the dashboard.\nMedian error ~74% — always verify before presenting to management."),
        (BLUE_LIGHT,  "Added date is preserved",
         "Every game remembers when it first appeared in the output.\nRe-running the pipeline never resets historical dates."),
        (GREEN,       "Model accuracy is honest",
         "IAA ~74% error  |  Blend ~33%  |  IAP High ~30%\nThis is a screening tool, not a valuation model. It surfaces candidates."),
        (AMBER,       "Notes & review tracking",
         "Checkboxes and notes per game are saved in the browser (localStorage).\nSummary card shows % reviewed and how many games have notes."),
    ]

    for i, (color, title, desc) in enumerate(items):
        row, col = divmod(i, 2)
        bx = 0.5 + col * 6.3
        by = 1.6 + row * 1.82
        card(s, bx, by, 5.9, 1.65, NAVY_CARD)
        accent_bar(s, bx, by, 1.65, color)
        txt(s, title, bx+0.2, by+0.12, 5.4, 0.38, size=14, bold=True, color=WHITE)
        txt(s, desc,  bx+0.2, by+0.58, 5.4, 0.95, size=11, color=OFF_WHITE)

    slide_number(s, n, total)

def slide_summary(prs, n, total):
    s = blank(prs); bg(s, NAVY)
    rect(s, 0, 0, 0.35, 7.5, BLUE)

    txt(s, "What we built", 0.8, 1.0, 11, 0.7,
        size=42, bold=True, color=WHITE)
    txt(s, "A fully automated, calibrated, and presentation-ready acquisition pipeline.",
        0.8, 1.85, 11, 0.5, size=18, color=BLUE_LIGHT)
    rect(s, 0.8, 2.5, 10, 0.04, BLUE)

    checkpoints = [
        "Profit model calibrated against 30 real games with 2025 actuals",
        "80+ publisher exclusions + auto-exclusion rule",
        "4-tab dashboard: Growing / Stable / Declining / Returning",
        "Cross-platform dedup: Google Play + App Store in one output",
        "Date tracking, notes, review progress — all persisted",
        "Upload button live today  ·  Execute + Salesforce coming next",
    ]
    for k, point in enumerate(checkpoints):
        by = 2.7 + k * 0.62
        rect(s, 0.8, by + 0.1, 0.28, 0.28, BLUE)
        txt(s, "✓", 0.8, by + 0.06, 0.28, 0.3, size=13, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, point, 1.25, by, 10.5, 0.5, size=15, color=OFF_WHITE)

    slide_number(s, n, total)

# ── Build ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    total = 10

    slide_title(prs)
    slide_intro(prs, 2, total)
    slide_model(prs, 3, total)
    slide_anchors(prs, 4, total)
    slide_exclusions(prs, 5, total)
    slide_tabs(prs, 6, total)
    slide_automation(prs, 7, total)
    slide_salesforce(prs, 8, total)
    slide_extras(prs, 9, total)
    slide_summary(prs, 10, total)

    out = r"C:\Rotem\AI Workshop\Publishing Dashboard\Game_Lead_Finder_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
